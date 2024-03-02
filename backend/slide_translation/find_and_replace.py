import os
from pptx import Presentation
from .translator import Translator

class FindAndReplace:
    def __init__(self):
        self.translator = Translator()

    def replace_text_in_slide(self, language, slide):
        """
        Replaces text in a slide with its corresponding translation.
        'replacements' is a dictionary where keys are original texts and values are their translations.
        """
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    original_text = run.text
                    translated_text = self.translator.translate(original_text, language)
                    run.text = translated_text

    def replace_text_in_presentation(self, file_name, language, presentation_path):
        """
        Loads a presentation, replaces text on each slide, and saves the presentation.
        """
        prs = Presentation(presentation_path)
        for slide in prs.slides:
            self.replace_text_in_slide(language, slide)

        if not os.path.exists('translated_presentations'):
            os.makedirs('translated_presentations')
        prs.save(f'translated_presentations/{file_name}_translated.pptx')
